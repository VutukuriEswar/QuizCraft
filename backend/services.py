import nltk
import os
import json
import re
import random
import time
import logging
import requests
from collections import Counter
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords, wordnet
from nltk.tag import pos_tag
from nltk.stem import WordNetLemmatizer

_TESSERACT_AVAILABLE = False
_PYMUPDF_AVAILABLE = False
_PIL_AVAILABLE = False
_DOCX_AVAILABLE = False
_PPTX_AVAILABLE = False

try:
    import pytesseract
    _TESSERACT_AVAILABLE = True
except:
    pass

try:
    import fitz
    _PYMUPDF_AVAILABLE = True
except:
    pass

try:
    from PIL import Image
    import io
    _PIL_AVAILABLE = True
except:
    pass

try:
    from docx import Document
    _DOCX_AVAILABLE = True
except:
    pass

try:
    from pptx import Presentation
    _PPTX_AVAILABLE = True
except:
    pass

STOP_WORDS = set(stopwords.words('english'))
LEMMA = WordNetLemmatizer()

EXTRA_STOPS = {
    'also', 'however', 'therefore', 'thus', 'hence', 'moreover',
    'furthermore', 'additionally', 'example', 'like', 'thing',
    'things', 'way', 'ways', 'use', 'used', 'using', 'make',
    'made', 'many', 'much', 'often', 'well', 'may', 'might',
    'can', 'could', 'would', 'should', 'one', 'two', 'first',
    'second', 'third', 'new', 'old', 'another', 'even', 'still',
    'since', 'yet', 'although', 'though', 'whether', 'rather',
    'either', 'neither', 'every', 'each', 'per', 'via', 'vs',
}

ALL_STOPS = STOP_WORDS | EXTRA_STOPS


class FileProcessor:

    @staticmethod
    def extract_text(file):
        if not file or not file.filename:
            return "", "No file provided"
        filename = file.filename.lower()
        try:
            if filename.endswith('.pdf'):
                return FileProcessor._extract_pdf(file)
            elif filename.endswith('.docx'):
                return FileProcessor._extract_docx(file)
            elif filename.endswith('.pptx'):
                return FileProcessor._extract_pptx(file)
            elif filename.endswith('.txt'):
                return FileProcessor._extract_txt(file)
            else:
                return "", "Unsupported file format. Use PDF, DOCX, PPTX, or TXT."
        except Exception as e:
            return "", f"Failed to extract text: {str(e)}"

    @staticmethod
    def _extract_pdf(file):
        if not _PYMUPDF_AVAILABLE:
            return "", "PDF processing unavailable. Install PyMuPDF."
        pdf_bytes = file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
            if _PIL_AVAILABLE and _TESSERACT_AVAILABLE:
                for img in page.get_images(full=True):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        img_obj = Image.open(io.BytesIO(image_bytes))
                        ocr_text = pytesseract.image_to_string(img_obj)
                        if ocr_text.strip():
                            text_parts.append(ocr_text)
                    except:
                        pass
        doc.close()
        full_text = "\n".join(text_parts).strip()
        if not full_text:
            return "", "Could not extract text from PDF"
        return full_text, None

    @staticmethod
    def _extract_docx(file):
        if not _DOCX_AVAILABLE:
            return "", "DOCX processing unavailable. Install python-docx."
        doc = Document(file)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                tables_text.append(" | ".join(row_text))
        full_text = "\n".join(paragraphs + tables_text).strip()
        if not full_text:
            return "", "Could not extract text from DOCX"
        return full_text, None

    @staticmethod
    def _extract_pptx(file):
        if not _PPTX_AVAILABLE:
            return "", "PPTX processing unavailable. Install python-pptx."
        prs = Presentation(file)
        text_parts = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(row_text))
            if slide_texts:
                text_parts.append("\n".join(slide_texts))
        full_text = "\n\n".join(text_parts).strip()
        if not full_text:
            return "", "Could not extract text from PPTX"
        return full_text, None

    @staticmethod
    def _extract_txt(file):
        content = file.read().decode('utf-8', errors='replace')
        return content.strip(), None


class TextPreprocessor:

    @staticmethod
    def preprocess(text):
        code_blocks = re.findall(r'```[\w]*\n(.*?)```', text, re.DOTALL)
        cleaned = re.sub(r'```[\w]*\n.*?```', '<<CODE_BLOCK>>', text, flags=re.DOTALL)
        cleaned = re.sub(r'```[^`\n]*```', '<<CODE_BLOCK>>', cleaned)
        cleaned = re.sub(r'[#*_`~>\[\]()]+', ' ', cleaned)
        cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
        cleaned = re.sub(r'[ \t]+', ' ', cleaned)
        return cleaned.strip(), code_blocks

    @staticmethod
    def get_sentences(text):
        cleaned, code_blocks = TextPreprocessor.preprocess(text)
        try:
            sentences = sent_tokenize(cleaned)
        except:
            sentences = re.split(r'(?<=[.!?])\s+', cleaned)
        result = []
        for sent in sentences:
            sent = sent.strip()
            if '<<CODE_BLOCK>>' in sent:
                parts = sent.split('<<CODE_BLOCK>>')
                for part in parts:
                    part = part.strip()
                    if part and len(part) > 10:
                        result.append(part)
            elif len(sent) > 15:
                result.append(sent)
        return result

    @staticmethod
    def get_code_blocks(text):
        blocks = re.findall(r'```[\w]*\n(.*?)```', text, re.DOTALL)
        inline = re.findall(r'```([^`\n]+)```', text)
        blocks.extend(inline)
        return blocks

    @staticmethod
    def get_code_context(text, code_block):
        idx = text.find(code_block[:50])
        if idx == -1:
            return ""
        start = max(0, idx - 300)
        end = min(len(text), idx + len(code_block) + 100)
        context = text[start:end]
        context = re.sub(r'```[\w]*\n.*?```', '', context, flags=re.DOTALL)
        return context.strip()


class ConceptExtractor:

    @staticmethod
    def extract(text):
        if not text or len(text.strip()) < 20:
            return []
        words = word_tokenize(text.lower())
        filtered = [w for w in words if w.isalpha() and w not in ALL_STOPS and len(w) > 2]
        if not filtered:
            return []
        freq = Counter(filtered)
        lemmatized = [LEMMA.lemmatize(w) for w in filtered]
        lemma_freq = Counter(lemmatized)
        tagged = pos_tag(word_tokenize(text))
        noun_phrases = ConceptExtractor._extract_noun_phrases(tagged)
        technical_terms = ConceptExtractor._find_technical_terms(text)
        concepts = []
        for word, count in lemma_freq.most_common(30):
            if count >= 2:
                concepts.append(word)
        seen_lower = {c.lower() for c in concepts}
        for phrase in noun_phrases:
            if phrase.lower() not in seen_lower and len(phrase) <= 40:
                concepts.append(phrase)
                seen_lower.add(phrase.lower())
        for term in technical_terms:
            if term.lower() not in seen_lower and 2 < len(term) <= 40:
                concepts.append(term)
                seen_lower.add(term.lower())
        return concepts[:50]

    @staticmethod
    def _extract_noun_phrases(tagged):
        phrases = []
        current = []
        for word, tag in tagged:
            if tag.startswith('NN'):
                current.append(word)
            elif tag.startswith('JJ') and current:
                current.append(word)
            else:
                if len(current) >= 2:
                    phrase = ' '.join(current)
                    if len(phrase) > 3 and phrase.lower() not in ALL_STOPS:
                        phrases.append(phrase)
                elif len(current) == 1 and current[0][0].isupper() and len(current[0]) > 3:
                    phrases.append(current[0])
                current = []
        if len(current) >= 2:
            phrase = ' '.join(current)
            if len(phrase) > 3:
                phrases.append(phrase)
        elif len(current) == 1 and current[0][0].isupper() and len(current[0]) > 3:
            phrases.append(current[0])
        return phrases

    @staticmethod
    def _find_technical_terms(text):
        terms = set()
        patterns = [
            r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+\b',
            r'\b[a-z]+_[a-z_]+\b',
            r'\b[A-Z]{2,}\b',
            r'\b\w+(?:::\w+)+\b',
            r'\b\w+\.\w+(?:\.\w+)*\b',
        ]
        for pattern in patterns:
            matches = re.findall(pattern, text)
            terms.update(matches)
        for match in re.findall(r'["\']([^"\']+)["\']', text):
            if 2 < len(match) <= 40:
                terms.add(match)
        for match in re.findall(r'`([^`]+)`', text):
            if 2 < len(match) <= 40:
                terms.add(match)
        return list(terms)


class DistractorGenerator:

    @staticmethod
    def get_distractors(target, concepts, full_text, count=3):
        target_lower = target.lower()
        candidates = []
        tagged_all = pos_tag(word_tokenize(full_text))
        target_tag = None
        for w, t in tagged_all:
            if w.lower() == target_lower:
                target_tag = t
                break
        if target_tag:
            for w, t in tagged_all:
                if t == target_tag and w.lower() != target_lower and w.lower() not in ALL_STOPS and len(w) > 2:
                    if w not in candidates:
                        candidates.append(w)
        for concept in concepts:
            if concept.lower() != target_lower and concept not in candidates:
                if 2 < len(concept) <= 30:
                    candidates.append(concept)
        synsets = wordnet.synsets(target)
        for synset in synsets[:3]:
            for lemma in synset.lemmas():
                name = lemma.name().replace('_', ' ')
                if name.lower() != target_lower and name not in candidates:
                    candidates.append(name)
        for synset in synsets[:2]:
            for hypernym in synset.hypernyms():
                for lemma in hypernym.lemmas():
                    name = lemma.name().replace('_', ' ')
                    if name.lower() != target_lower and name not in candidates and len(name.split()) <= 3:
                        candidates.append(name)
        for synset in synsets[:2]:
            for hyponym in synset.hyponyms()[:5]:
                for lemma in hyponym.lemmas():
                    name = lemma.name().replace('_', ' ')
                    if name.lower() != target_lower and name not in candidates and len(name.split()) <= 3:
                        candidates.append(name)
        text_words = set(w for w in word_tokenize(full_text) if w.isalpha() and len(w) > 3 and w.lower() not in ALL_STOPS)
        for word in text_words:
            if word.lower() not in [c.lower() for c in candidates] and word.lower() != target_lower:
                candidates.append(word)
        filtered = []
        for c in candidates:
            if c.lower() != target_lower:
                similarity = DistractorGenerator._string_similarity(target, c)
                if 0.1 < similarity < 0.8:
                    filtered.append(c)
                elif similarity >= 0.8:
                    continue
                else:
                    filtered.append(c)
        result = []
        seen = set()
        for c in filtered:
            if c.lower() not in seen:
                seen.add(c.lower())
                result.append(c)
            if len(result) >= count:
                break
        if len(result) < count:
            for c in candidates:
                if c.lower() not in seen and c.lower() != target_lower:
                    seen.add(c.lower())
                    result.append(c)
                if len(result) >= count:
                    break
        return result[:count]

    @staticmethod
    def _string_similarity(a, b):
        a, b = a.lower(), b.lower()
        if not a or not b:
            return 0.0
        set_a, set_b = set(a), set(b)
        intersection = set_a & set_b
        union = set_a | set_b
        if not union:
            return 0.0
        return len(intersection) / len(union)


class SentenceClassifier:

    @staticmethod
    def classify(sentence):
        lower = sentence.lower()
        categories = []
        if any(p in lower for p in ['is defined as', 'refers to', 'means', 'is known as', 'is called', 'definition']):
            categories.append('definition')
        if any(p in lower for p in ['used for', 'purpose', 'in order to', 'designed to', 'aims to', 'serves to']):
            categories.append('purpose')
        if any(p in lower for p in ['because', 'due to', 'results in', 'leads to', 'causes', 'affects', 'depends on']):
            categories.append('causal')
        if any(p in lower for p in ['first', 'then', 'next', 'finally', 'step', 'process', 'followed by']):
            categories.append('process')
        if any(p in lower for p in ['compared to', 'unlike', 'similar to', 'difference', 'whereas', 'while']):
            categories.append('comparison')
        if any(p in lower for p in ['example', 'for instance', 'such as', 'including', 'like']):
            categories.append('example')
        if any(p in lower for p in ['contains', 'consists of', 'comprises', 'includes', 'made up of']):
            categories.append('compositional')
        if re.search(r'\d+', sentence):
            categories.append('numerical')
        if not categories:
            categories.append('general')
        return categories


class NLPQuestionGenerator:

    @staticmethod
    def generate(text, concepts, config_blocks):
        questions = []
        sentences = TextPreprocessor.get_sentences(text)
        code_blocks = TextPreprocessor.get_code_blocks(text)
        ranked = NLPQuestionGenerator._rank_sentences(sentences, concepts)
        for block in config_blocks:
            q_type = block.get('type', 'mcq')
            count = block.get('count', 5)
            marks = block.get('marks', 1)
            if q_type == 'mcq':
                generated = NLPQuestionGenerator._generate_mcq(ranked, concepts, text, count, marks)
            elif q_type == 'msq':
                generated = NLPQuestionGenerator._generate_msq(ranked, concepts, text, count, marks)
            elif q_type == 'true_false':
                generated = NLPQuestionGenerator._generate_tf(ranked, concepts, count, marks)
            else:
                generated = NLPQuestionGenerator._generate_mcq(ranked, concepts, text, count, marks)
            questions.extend(generated)
        if code_blocks:
            mcq_blocks = [b for b in config_blocks if b.get('type') == 'mcq']
            if mcq_blocks:
                total_mcq = sum(b.get('count', 0) for b in mcq_blocks)
                mcq_marks = mcq_blocks[0].get('marks', 1)
                code_qs = NLPQuestionGenerator._generate_code_questions(code_blocks, text, total_mcq, mcq_marks, concepts)
                questions.extend(code_qs)
        questions = NLPQuestionGenerator._deduplicate(questions)
        questions = NLPQuestionGenerator._validate(questions)
        return questions

    @staticmethod
    def _rank_sentences(sentences, concepts):
        concept_lower = [c.lower() for c in concepts]
        scored = []
        for i, sent in enumerate(sentences):
            sent_lower = sent.lower()
            words = word_tokenize(sent_lower)
            content_words = [w for w in words if w.isalpha() and w not in ALL_STOPS]
            if len(content_words) < 4:
                continue
            concept_hits = sum(1 for c in concept_lower if c in sent_lower)
            if concept_hits == 0:
                continue
            unique_ratio = len(set(content_words)) / max(len(content_words), 1)
            length_score = min(len(content_words) / 15, 1.0)
            position_score = 1.0 if i < max(len(sentences) * 0.3, 5) else 0.7
            has_verb = any(tag.startswith('VB') for _, tag in pos_tag(word_tokenize(sent)))
            verb_score = 1.0 if has_verb else 0.5
            categories = SentenceClassifier.classify(sent)
            category_bonus = 1.5 if 'definition' in categories else (1.2 if 'purpose' in categories else 1.0)
            score = (concept_hits * 3) + (unique_ratio * 2) + (length_score * 1.5) + (position_score) + (verb_score * 2) * category_bonus
            scored.append((sent, score, categories))
        scored.sort(key=lambda x: x[1], reverse=True)
        return [(s[0], s[2]) for s in scored]

    @staticmethod
    def _generate_mcq(sentences, concepts, full_text, count, marks):
        questions = []
        used = set()
        definition_sents = [(s, c) for s, c in sentences if 'definition' in c]
        purpose_sents = [(s, c) for s, c in sentences if 'purpose' in c]
        causal_sents = [(s, c) for s, c in sentences if 'causal' in c]
        general_sents = [(s, c) for s, c in sentences if 'general' in c]
        ordered = definition_sents + purpose_sents + causal_sents + general_sents
        for sent, categories in ordered:
            if len(questions) >= count or sent in used:
                continue
            result = None
            if 'definition' in categories:
                result = NLPQuestionGenerator._mcq_definition(sent, concepts, full_text)
            elif 'purpose' in categories:
                result = NLPQuestionGenerator._mcq_purpose(sent, concepts, full_text)
            elif 'causal' in categories:
                result = NLPQuestionGenerator._mcq_causal(sent, concepts, full_text)
            elif 'compositional' in categories:
                result = NLPQuestionGenerator._mcq_compositional(sent, concepts, full_text)
            elif 'numerical' in categories:
                result = NLPQuestionGenerator._mcq_numerical(sent, concepts, full_text)
            if not result:
                result = NLPQuestionGenerator._mcq_general(sent, concepts, full_text)
            if result and len(result.get('options', [])) >= 4:
                used.add(sent)
                questions.append({
                    "question_text": result["question"],
                    "options": result["options"],
                    "answer": result["answer"],
                    "type": "mcq",
                    "marks": marks
                })
        return questions[:count]

    @staticmethod
    def _get_key_nouns(sentence):
        tagged = pos_tag(word_tokenize(sentence))
        nouns = []
        for word, tag in tagged:
            if tag.startswith('NN') and word.lower() not in ALL_STOPS and len(word) > 2:
                nouns.append((word, tag))
        return nouns

    @staticmethod
    def _mcq_definition(sent, concepts, full_text):
        nouns = NLPQuestionGenerator._get_key_nouns(sent)
        if not nouns:
            return None
        target_word, target_tag = max(nouns, key=lambda x: len(x[0]))
        distractors = DistractorGenerator.get_distractors(target_word, concepts, full_text, 3)
        if len(distractors) < 3:
            return None
        templates = [
            f"Which of the following correctly identifies what '{target_word}' refers to in the given context?",
            f"Based on the description, '{target_word}' is best understood as:",
            f"What is the meaning of '{target_word}' as described in the text?",
            f"The term '{target_word}' in this context represents which of the following?",
        ]
        if target_tag == 'NNS':
            templates = [
                f"Which option best describes the '{target_word}' mentioned in the text?",
                f"What are the '{target_word}' according to the given description?",
            ]
        options = [target_word] + distractors[:3]
        random.shuffle(options)
        return {"question": random.choice(templates), "options": options, "answer": target_word}

    @staticmethod
    def _mcq_purpose(sent, concepts, full_text):
        nouns = NLPQuestionGenerator._get_key_nouns(sent)
        if not nouns:
            return None
        target_word = max(nouns, key=lambda x: len(x[0]))[0]
        distractors = DistractorGenerator.get_distractors(target_word, concepts, full_text, 3)
        if len(distractors) < 3:
            return None
        templates = [
            f"What is the primary purpose of '{target_word}' as described in the text?",
            f"Why is '{target_word}' used according to the given information?",
            f"The main function of '{target_word}' in this context is to:",
            f"Which statement best explains why '{target_word}' is necessary?",
        ]
        options = [target_word] + distractors[:3]
        random.shuffle(options)
        return {"question": random.choice(templates), "options": options, "answer": target_word}

    @staticmethod
    def _mcq_causal(sent, concepts, full_text):
        nouns = NLPQuestionGenerator._get_key_nouns(sent)
        if len(nouns) < 2:
            return None
        cause_word = nouns[0][0]
        effect_word = nouns[-1][0]
        distractors = DistractorGenerator.get_distractors(effect_word, concepts, full_text, 3)
        if len(distractors) < 3:
            return None
        templates = [
            f"According to the text, what is the effect of '{cause_word}'?",
            f"What results from '{cause_word}' as described in the passage?",
            f"Which of the following is a consequence of '{cause_word}'?",
            f"How does '{cause_word}' influence the outcome described?",
        ]
        options = [effect_word] + distractors[:3]
        random.shuffle(options)
        return {"question": random.choice(templates), "options": options, "answer": effect_word}

    @staticmethod
    def _mcq_compositional(sent, concepts, full_text):
        nouns = NLPQuestionGenerator._get_key_nouns(sent)
        if not nouns:
            return None
        target_word = max(nouns, key=lambda x: len(x[0]))[0]
        distractors = DistractorGenerator.get_distractors(target_word, concepts, full_text, 3)
        if len(distractors) < 3:
            return None
        templates = [
            f"Which of the following is stated to be a component or part of '{target_word}'?",
            f"What does '{target_word}' consist of according to the text?",
            f"Which item is included in '{target_word}' as described?",
        ]
        options = [target_word] + distractors[:3]
        random.shuffle(options)
        return {"question": random.choice(templates), "options": options, "answer": target_word}

    @staticmethod
    def _mcq_numerical(sent, concepts, full_text):
        numbers = re.findall(r'\b\d+(?:\.\d+)?\b', sent)
        if not numbers:
            return NLPQuestionGenerator._mcq_general(sent, concepts, full_text)
        target_num = random.choice(numbers)
        try:
            base = float(target_num)
            if base == int(base):
                base = int(base)
                variants = [base + random.randint(1, 5) * random.choice([-1, 1]) for _ in range(3)]
            else:
                variants = [round(base + random.uniform(0.5, 3.0) * random.choice([-1, 1]), 1) for _ in range(3)]
        except:
            return NLPQuestionGenerator._mcq_general(sent, concepts, full_text)
        distractors = [str(v) for v in variants]
        context = re.sub(r'\b\d+(?:\.\d+)?\b', '____', sent, count=1)
        if len(context) < 20:
            return NLPQuestionGenerator._mcq_general(sent, concepts, full_text)
        templates = [
            f"What value completes the following statement? \"{context}\"",
            f"Fill in the blank: {context}",
            f"Based on the text, the correct value in \"{context}\" is:",
        ]
        options = [target_num] + distractors
        random.shuffle(options)
        return {"question": random.choice(templates), "options": options, "answer": target_num}

    @staticmethod
    def _mcq_general(sent, concepts, full_text):
        nouns = NLPQuestionGenerator._get_key_nouns(sent)
        if not nouns:
            return None
        target_word, target_tag = max(nouns, key=lambda x: len(x[0]))
        distractors = DistractorGenerator.get_distractors(target_word, concepts, full_text, 3)
        if len(distractors) < 3:
            return None
        templates = [
            f"Which of the following best describes the role of '{target_word}' in the given context?",
            f"In the described scenario, what is most accurately associated with '{target_word}'?",
            f"Based on the information provided, which option correctly relates to '{target_word}'?",
            f"What does the text indicate about '{target_word}'?",
        ]
        options = [target_word] + distractors[:3]
        random.shuffle(options)
        return {"question": random.choice(templates), "options": options, "answer": target_word}

    @staticmethod
    def _generate_msq(sentences, concepts, full_text, count, marks):
        questions = []
        used_sents = set()
        list_candidates = []
        for sent, categories in sentences:
            if any(c in categories for c in ['example', 'compositional']):
                items = NLPQuestionGenerator._extract_list_items(sent)
                if len(items) >= 3:
                    list_candidates.append((sent, items))
        for sent, items in list_candidates:
            if len(questions) >= count or sent in used_sents:
                continue
            correct_count = random.randint(2, min(len(items), 4))
            correct_items = random.sample(items, correct_count)
            other_concepts = [c for c in concepts if c.lower() not in [i.lower() for i in items] and 2 < len(c) <= 30]
            distractor_count = min(3, len(other_concepts))
            if distractor_count < 2:
                continue
            distractor_items = random.sample(other_concepts, distractor_count)
            all_options = correct_items + distractor_items
            random.shuffle(all_options)
            context = re.sub(r'(?:such as|including|like|for example|e\.g\.)[^.]*', '...', sent, flags=re.IGNORECASE)
            context = context.strip().rstrip(',')
            if len(context) < 15:
                context = sent[:80]
            templates = [
                f"Based on the context: \"{context}\", which of the following are correct?",
                f"Which of the following are associated with the described concept?",
                f"Select all that apply according to the given information.",
            ]
            used_sents.add(sent)
            questions.append({
                "question_text": random.choice(templates),
                "options": all_options,
                "answer": correct_items,
                "type": "msq",
                "marks": marks
            })
        if len(questions) < count:
            for sent, categories in sentences:
                if len(questions) >= count or sent in used_sents:
                    continue
                nouns = NLPQuestionGenerator._get_key_nouns(sent)
                if len(nouns) >= 5:
                    correct = random.sample(nouns, random.randint(2, 3))
                    remaining = [n for n in nouns if n not in correct]
                    if len(remaining) >= 2:
                        distractors = random.sample(remaining, 2)
                        all_opts = [n[0] if isinstance(n, tuple) else n for n in correct + distractors]
                        random.shuffle(all_opts)
                        used_sents.add(sent)
                        questions.append({
                            "question_text": "Which of the following concepts are directly mentioned or implied in the given context?",
                            "options": all_opts,
                            "answer": [n[0] if isinstance(n, tuple) else n for n in correct],
                            "type": "msq",
                            "marks": marks
                        })
        return questions[:count]

    @staticmethod
    def _extract_list_items(sentence):
        items = []
        patterns = [
            r'such as\s+((?:[^,.]+(?:,\s*)?)+)',
            r'including\s+((?:[^,.]+(?:,\s*)?)+)',
            r'(?:for example|e\.g\.)\s*((?:[^,.]+(?:,\s*)?)+)',
        ]
        for pattern in patterns:
            match = re.search(pattern, sentence, re.IGNORECASE)
            if match:
                list_str = match.group(1)
                items = [item.strip() for item in re.split(r',\s*(?:and\s*)?', list_str) if item.strip() and 1 < len(item.strip()) <= 40]
                if len(items) >= 3:
                    return items
        clauses = re.split(r',\s+', sentence)
        if len(clauses) >= 3:
            for clause in clauses:
                clause = clause.strip()
                if 1 < len(clause) < 40:
                    words = word_tokenize(clause)
                    if any(t.startswith('NN') for _, t in pos_tag(words)):
                        items.append(clause)
        if len(items) >= 3:
            return items
        numbered = re.findall(r'(?:\d+[\.\)]\s*)([^.\n]+)', sentence)
        if len(numbered) >= 3:
            return [n.strip() for n in numbered if len(n.strip()) > 1]
        return items

    @staticmethod
    def _generate_tf(sentences, concepts, count, marks):
        questions = []
        used = set()
        factual_patterns = [
            r'\bis\b', r'\bare\b', r'\bwas\b', r'\bwere\b',
            r'\bhas\b', r'\bhave\b', r'\bhad\b',
            r'\bcan\b', r'\bcannot\b', r'\bwill\b',
            r'\bdoes\b', r'\brequires?\b', r'\bprovides?\b',
            r'\bconsists?\b', r'\bcontains?\b',
            r'\bdefines?\b', r'\bdescribes?\b',
            r'\balways\b', r'\bnever\b', r'\bonly\b',
        ]
        factual = []
        for sent, categories in sentences:
            if any(re.search(p, sent, re.IGNORECASE) for p in factual_patterns):
                words = word_tokenize(sent)
                if 5 <= len(words) <= 45:
                    factual.append(sent)
        for sent in factual:
            if len(questions) >= count or sent in used:
                continue
            used.add(sent)
            is_true = random.choice([True, False])
            if is_true:
                statement = sent
                answer = "True"
            else:
                statement = NLPQuestionGenerator._make_false(sent, concepts)
                if statement == sent:
                    statement = sent
                    answer = "True"
                else:
                    answer = "False"
            questions.append({
                "question_text": statement,
                "options": ["True", "False"],
                "answer": answer,
                "type": "true_false",
                "marks": marks
            })
        if len(questions) < count:
            for sent, _ in sentences:
                if len(questions) >= count or sent in used:
                    continue
                words = word_tokenize(sent)
                if 5 <= len(words) <= 45:
                    used.add(sent)
                    questions.append({
                        "question_text": sent,
                        "options": ["True", "False"],
                        "answer": "True",
                        "type": "true_false",
                        "marks": marks
                    })
        return questions[:count]

    @staticmethod
    def _make_false(sentence, concepts):
        words = word_tokenize(sentence)
        tagged = pos_tag(words)
        nouns = [(i, w) for i, (w, t) in enumerate(tagged) if t.startswith('NN') and w.lower() not in ALL_STOPS and len(w) > 2]
        if len(nouns) >= 2:
            i1, w1 = nouns[0]
            i2, w2 = nouns[-1]
            if w1.lower() != w2.lower():
                new_words = words[:]
                new_words[i1] = w2
                new_words[i2] = w1
                return ' '.join(new_words)
        negation_map = {
            'is': 'is not', 'are': 'are not',
            'was': 'was not', 'were': 'were not',
            'has': 'does not have', 'have': 'do not have',
            'can': 'cannot', 'will': 'will not',
            'does': 'does not', 'always': 'never',
            'never': 'always', 'all': 'not all',
        }
        for i, word in enumerate(words):
            if word.lower() in negation_map:
                new_words = words[:]
                new_words[i] = negation_map[word.lower()]
                return ' '.join(new_words)
        numbers = [(i, w) for i, w in enumerate(words) if w.isdigit()]
        if numbers:
            idx, num = random.choice(numbers)
            new_num = str(int(num) + random.choice([1, 2, 3, 5, 10]))
            new_words = words[:]
            new_words[idx] = new_num
            return ' '.join(new_words)
        adjectives = [(i, w) for i, (w, t) in enumerate(tagged) if t.startswith('JJ') and len(w) > 3 and w.lower() not in ALL_STOPS]
        if adjectives:
            idx, adj = random.choice(adjectives)
            antonyms = []
            for synset in wordnet.synsets(adj):
                for lemma in synset.lemmas():
                    if lemma.antonyms():
                        antonyms.extend(a.name().replace('_', ' ') for a in lemma.antonyms())
            if antonyms:
                new_words = words[:]
                new_words[idx] = antonyms[0]
                return ' '.join(new_words)
        return sentence

    @staticmethod
    def _generate_code_questions(code_blocks, full_text, total_mcq, marks, concepts):
        questions = []
        code_count = min(len(code_blocks), max(1, total_mcq // 4))
        for code in code_blocks[:code_count]:
            result = NLPQuestionGenerator._create_code_question(code, full_text, marks)
            if result:
                questions.append(result)
        return questions

    @staticmethod
    def _create_code_question(code, full_text, marks):
        lines = [l.strip() for l in code.strip().split('\n') if l.strip()]
        if len(lines) < 2:
            return None
        func_match = re.search(r'(?:def|function|func|public\s+\w+\s+|private\s+\w+\s+|static\s+\w+\s+)(\w+)', code)
        class_match = re.search(r'(?:class|struct|interface|trait)\s+(\w+)', code)
        variables = re.findall(r'(\w+)\s*=[^=]', code)
        variables = [v for v in variables if v not in ('if', 'else', 'for', 'while', 'return', 'print', 'self', 'None', 'True', 'False') and not v[0].isdigit()]
        imports = re.findall(r'(?:import|from)\s+([\w.]+)', code)
        loops = len(re.findall(r'\b(for|while|do)\b', code))
        conditionals = len(re.findall(r'\b(if|elif|else|switch|case)\b', code))
        return_matches = re.findall(r'return\s+(.+)', code)
        context = TextPreprocessor.get_code_context(full_text, code)
        question = None
        answer = None
        if func_match:
            func_name = func_match.group(1)
            if return_matches:
                ret_val = return_matches[0].strip()[:50]
                templates = [
                    f"What does the function '{func_name}' return in the given code?",
                    f"The return value of '{func_name}' in this code is:",
                    f"What will '{func_name}' output when called?",
                ]
                answer = ret_val
            else:
                templates = [
                    f"What is the primary purpose of the function '{func_name}'?",
                    f"The function '{func_name}' is designed to accomplish which task?",
                    f"Which of the following best describes what '{func_name}' does?",
                ]
                if context:
                    context_words = [w for w in word_tokenize(context.lower()) if w.isalpha() and w not in ALL_STOPS]
                    answer = ' '.join(context_words[:6]) if context_words else f"Implements {func_name} logic"
                else:
                    answer = f"Implements the logic of {func_name}"
            question = random.choice(templates)
        elif class_match:
            class_name = class_match.group(1)
            templates = [
                f"What does the class '{class_name}' represent in the given code?",
                f"The class '{class_name}' is designed to model which concept?",
                f"Which of the following best describes the purpose of '{class_name}'?",
            ]
            answer = f"Defines the structure and behavior of {class_name}"
            question = random.choice(templates)
        elif variables:
            var = variables[0]
            templates = [
                f"What does the variable '{var}' store in the given code?",
                f"In the provided code, what is the purpose of '{var}'?",
                f"The variable '{var}' is assigned what kind of value?",
            ]
            answer = f"Stores a value related to {var}"
            question = random.choice(templates)
        elif loops > 0:
            templates = [
                "What is the primary purpose of the loop in the given code?",
                "The loop construct in this code is used to accomplish which task?",
                "What does the iterative structure in this code achieve?",
            ]
            answer = "Iterates to perform repeated operations on data"
            question = random.choice(templates)
        elif conditionals > 0:
            templates = [
                "What condition is being checked in the given code?",
                "The conditional logic in this code controls which behavior?",
                "What decision is being made by the if/else block in this code?",
            ]
            answer = "Controls conditional execution based on a condition"
            question = random.choice(templates)
        elif imports:
            lib = imports[0]
            templates = [
                f"Why is '{lib}' imported in the given code?",
                f"What functionality does '{lib}' provide in this code?",
                f"The import of '{lib}' suggests this code is related to which area?",
            ]
            answer = f"Provides functionality from the {lib} library"
            question = random.choice(templates)
        else:
            templates = [
                "What is the overall purpose of the given code snippet?",
                "Which of the following best describes what the provided code accomplishes?",
                "The given code primarily serves which purpose?",
            ]
            answer = "Executes the described programming logic"
            question = random.choice(templates)
        distractors = [
            "Handles error handling and exception management",
            "Manages database connections and SQL queries",
            "Performs user input validation and sanitization",
            "Creates graphical user interface components",
            "Manages network socket communications",
            "Implements encryption and security protocols",
            "Parses and processes XML/JSON data structures",
            "Coordinates multi-threaded parallel execution",
            "Handles file system operations and I/O",
            "Implements caching and memory management",
            "Performs unit testing and assertions",
            "Manages API endpoint routing and middleware",
        ]
        random.shuffle(distractors)
        options = [answer] + distractors[:3]
        random.shuffle(options)
        display_code = '\n'.join(lines[:12])
        if len(lines) > 12:
            display_code += '\n  ...'
        return {
            "question_text": f"{question}\n\n```\n{display_code}\n```",
            "options": options,
            "answer": answer,
            "type": "mcq",
            "marks": marks
        }

    @staticmethod
    def _deduplicate(questions):
        seen = set()
        unique = []
        for q in questions:
            q_text = q.get('question_text', '').lower().strip()
            if q_text and q_text not in seen:
                seen.add(q_text)
                unique.append(q)
        return unique

    @staticmethod
    def _validate(questions):
        valid = []
        for q in questions:
            q_text = q.get('question_text', '')
            if not q_text or len(q_text) < 15:
                continue
            q_type = q.get('type', 'mcq')
            if q_type == 'mcq':
                options = q.get('options', [])
                answer = q.get('answer', '')
                if not options or len(options) < 2:
                    continue
                if answer not in options:
                    continue
                if len(set(str(o).lower() for o in options)) < len(options):
                    continue
            elif q_type == 'msq':
                options = q.get('options', [])
                answer = q.get('answer', '')
                if not options or len(options) < 3:
                    continue
                if not isinstance(answer, list) or len(answer) == 0:
                    continue
                if not all(a in options for a in answer):
                    continue
            elif q_type == 'true_false':
                answer = q.get('answer', '')
                if str(answer).strip() not in ['True', 'False']:
                    continue
            if len(q_text) > 1000:
                q['question_text'] = q_text[:1000]
            valid.append(q)
        return valid


class AIGenerator:

    @staticmethod
    def check_key_valid(api_key):
        try:
            response = requests.get(
                url="https://openrouter.ai/api/v1/key",
                headers={
                    "Authorization": f"Bearer {api_key}"
                },
                timeout=5
            )
            if response.status_code == 200:
                data = response.json()
                logging.info(f"OPENROUTER KEY DETAILS: {json.dumps(data, indent=2)}")
                return not data.get("error")
            else:
                logging.warning(f"OPENROUTER KEY CHECK FAILED: {response.status_code}")
                return False
        except requests.exceptions.ConnectionError:
            logging.warning("OPENROUTER UNREACHABLE: Could not connect to OpenRouter API. Falling back to NLP generation.")
            return False
        except requests.exceptions.RequestException as e:
            logging.warning(f"OPENROUTER API ERROR: Falling back to NLP generation. Error: {str(e)}")
            return False
        except Exception as e:
            logging.error(f"OPENROUTER KEY CHECK EXCEPTION: {str(e)}")
            return False

    @staticmethod
    def _clean_json_string(raw_string):
        result = []
        in_string = False
        escape_next = False
        for char in raw_string:
            if escape_next:
                result.append(char)
                escape_next = False
                continue
            if char == '\\':
                result.append(char)
                escape_next = True
                continue
            if char == '"':
                in_string = not in_string
                result.append(char)
                continue
            if in_string and char in '\n\r\t':
                result.append('\\n')
            else:
                result.append(char)
        return "".join(result)

    @staticmethod
    def generate_with_ai(input_text, concept_list, config_blocks):
        print(">>> [QUIZCRAFT DEBUG] ENTERED AI GENERATOR <<<")
        api_key = os.environ.get("OPENROUTER_API_KEY")
        if not api_key:
            print(">>> [QUIZCRAFT DEBUG] NO API KEY FOUND <<<")
            return None

        config_desc = []
        type_details = []
        for block in config_blocks:
            q_type = block.get('type', 'mcq')
            count = block.get('count', 5)
            marks = block.get('marks', 1)
            config_desc.append(f"{count} {q_type}")
            if q_type == 'mcq':
                type_details.append(f"- Generate exactly {count} MCQ questions. Each must have \"type\": \"mcq\", an \"options\" array with exactly 4 options, and \"answer\" must be the EXACT FULL TEXT of the correct option (copy it exactly from the options array). NEVER use letters like A, B, C, D as the answer. \"marks\": {marks}")
            elif q_type == 'msq':
                type_details.append(f"- Generate exactly {count} MSQ questions. Each must have \"type\": \"msq\", an \"options\" array with exactly 5 options, and \"answer\" MUST be a JSON ARRAY containing the EXACT FULL TEXT of each correct option (copy them exactly from the options array). The answer MUST always be an array with 2 or 3 items, NEVER a single string. NEVER use letters like A, B, C, D in the answer array. \"marks\": {marks}")
            elif q_type == 'true_false':
                type_details.append(f"- Generate exactly {count} True/False questions. Each must have \"type\": \"true_false\", \"options\": [\"True\", \"False\"], and \"answer\" must be either \"True\" or \"False\" (as a string, not an array). \"marks\": {marks}")
        config_str = ", ".join(config_desc)
        type_instructions = "\n".join(type_details)
        has_code = bool(re.search(r'```[\w]*\n', input_text) or re.search(r'(?:def |class |function |import |const |let |var |public |private )', input_text))
        code_instruction = ""
        if has_code:
            code_instruction = """

CRITICAL CODE INSTRUCTION: The input contains code. You may ONLY ask about the EXACT variables, functions, classes, and logic explicitly written in the provided code snippet. 
You are FORBIDDEN from using external programming examples (e.g., do NOT mention Fibonacci, Factorial, Sorting algorithms, Palindromes, or any other concept not explicitly written in the provided text/code).
"""
        system_prompt = """You are an expert educator and question designer.

========================
INPUT TEXT:
{{input_text}}

EXTRACTED CONCEPTS:
{{concept_list}}
========================

INSTRUCTIONS:

1. Deeply analyze the input text and understand:
   - Core concepts
   - Relationships between concepts
   - Explanations, examples, and logic present in the text

2. Generate questions such that:
   - Each question MUST be answerable using ONLY the given input text
   - Questions should test UNDERSTANDING, not memorization
   - A student who understands the concept in the text should be able to solve them
   - Avoid questions that require outside knowledge

3. Ensure diversity in thinking levels:
   - Conceptual understanding (WHY / HOW)
   - Application-based (apply concept to scenario)
   - Logical reasoning (connect multiple ideas from text)
   - Edge/tricky questions (test depth of understanding)
""" + code_instruction + """
4. ZERO HALLUCINATION RULES (STRICT ENFORCEMENT):
   - DO NOT mention any concept, term, algorithm, or proper noun that does not explicitly appear in the INPUT TEXT above.
   - DO NOT use your internal knowledge to fill in gaps.
   - If the text is about "Reinforcement Learning", ONLY ask about Reinforcement Learning as described in the text. Do NOT ask about generic machine learning concepts unless they are in the text.
   - If the text contains code, ONLY ask about the exact functions/variables written in that code. Do NOT substitute generic code examples (like Fibonacci, Factorial, etc).
   - If you cannot generate the requested number of questions without hallucinating outside terms, generate as many as you can without hallucinating, or return an empty array [].
   - VIOLATING THIS RULE WILL RESULT IN FAILURE.

5. EXAMPLE CONTEXT RULE (IMPORTANT):
   - When a question references a specific example, scenario, case study, or illustration from the input text, you MUST include a brief summary of that example's key details directly in the question text itself.
   - The question must be self-contained and answerable without the reader needing to recall the original example from memory.
   - WRONG: "In the motorbike rental company example, what assumption makes this suitable for Dynamic Programming?"
   - RIGHT: "In the example of a motorbike rental company that transfers bikes between locations based on random demand and return rates, what assumption makes this problem suitable for Dynamic Programming?"
   - Apply this to ANY specific example, scenario, or case study referenced from the input text.

6. For each question, also provide:
   - Correct answer
   - Short explanation (based ONLY on input text)

QUESTION FORMAT REQUIREMENTS:
""" + type_instructions + """

========================
CRITICAL ANSWER FORMAT RULES (FAILURE TO FOLLOW WILL REJECT YOUR OUTPUT):

FOR MCQ QUESTIONS:
- "answer" MUST be the EXACT FULL TEXT of the correct option
- "answer" MUST be a STRING, not an array
- NEVER put just a letter like "A", "B", "C", or "D" as the answer
- The answer string MUST match one of the strings in the "options" array EXACTLY (character-for-character)

EXAMPLE OF CORRECT MCQ:
{
  "question": "What is X?",
  "type": "mcq",
  "thinking": "conceptual",
  "concept": "X",
  "answer": "The environment's probability distributions are known",
  "options": ["The company can accurately predict customer behavior", "The environment's probability distributions are known", "The number of bikes is unlimited", "All states are equally likely"],
  "explanation": "The text states that DP requires known probability distributions."
}

FOR MSQ QUESTIONS:
- "answer" MUST be a JSON ARRAY of strings
- NEVER put a single string as the answer for MSQ
- Each string in the answer array MUST be the EXACT FULL TEXT of a correct option
- NEVER put letters like "A", "B", "C", "D", "E" in the answer array
- Each string in the answer array MUST match one of the strings in the "options" array EXACTLY

EXAMPLE OF CORRECT MSQ:
{
  "question": "Which of the following are characteristics of Monte Carlo methods?",
  "type": "msq",
  "thinking": "conceptual",
  "concept": "Monte Carlo",
  "answer": ["They learn from sample episodes", "They update values after complete episodes"],
  "options": ["They learn from sample episodes", "They require knowledge of MDP transitions", "They can only be applied to episodic problems", "They use bootstrapping for value estimation", "They update values after complete episodes"],
  "explanation": "Monte Carlo methods learn from complete episodes without bootstrapping."
}

FOR TRUE/FALSE QUESTIONS:
- "answer" MUST be either the string "True" or the string "False"
- "answer" MUST be a STRING, not an array
- NEVER put anything other than exactly "True" or "False"

========================
OUTPUT FORMAT (STRICT JSON ARRAY — no markdown, no code fences, raw JSON only):

[
  {
    "question": "...",
    "type": "mcq",
    "thinking": "conceptual",
    "concept": "...",
    "answer": "EXACT TEXT OF CORRECT OPTION FROM OPTIONS ARRAY",
    "options": ["Option 1 full text", "Option 2 full text", "Option 3 full text", "Option 4 full text"],
    "explanation": "..."
  }
]
========================"""

        user_prompt = f"""Generate exactly these questions from the text: {config_str}

REMEMBER:
- For MCQ: "answer" must be the FULL TEXT of the correct option, NEVER a letter like A/B/C/D
- For MSQ: "answer" must be a JSON ARRAY of FULL TEXTS of correct options, NEVER a single string and NEVER letters
- For True/False: "answer" must be exactly "True" or "False"

Ensure each question tests deep understanding. Mix conceptual, application, reasoning, and tricky question types in the "thinking" field.
{"Include questions about the EXACT code found in the text." if has_code else ""}
Return ONLY a raw JSON array. No markdown fences, no explanation, just the JSON."""

        formatted_system = system_prompt.replace("{{input_text}}", input_text).replace("{{concept_list}}", json.dumps(concept_list[:40]))

        max_retries = 15
        for attempt in range(max_retries):
            try:
                print(f">>> [QUIZCRAFT DEBUG] SENDING REQUEST TO OPENROUTER (Attempt {attempt + 1}/{max_retries}) <<<")
                response = requests.post(
                    url="https://openrouter.ai/api/v1/chat/completions",
                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "Content-Type": "application/json",
                        "HTTP-Referer": "http://localhost:3000",
                        "X-OpenRouter-Title": "QuizCraft"
                    },
                    data=json.dumps({
                        "model": "z-ai/glm-4.5-air:free",
                        "messages": [
                            {"role": "system", "content": formatted_system},
                            {"role": "user", "content": user_prompt}
                        ],
                        "temperature": 0.7,
                        "max_tokens": 6000
                    })
                )
                response.raise_for_status()
                result = response.json()
                content = result['choices'][0]['message'].get('content', '')

                print(f">>> [QUIZCRAFT DEBUG] RECEIVED RESPONSE. LENGTH: {len(content)} <<<")
                if not content:
                    logging.error("OPENROUTER FALLBACK: AI returned completely empty content.")
                    return None
                if len(content) < 50:
                    logging.error(f"OPENROUTER FALLBACK: AI returned suspiciously short content: {content}")
                    return None

                content = re.sub(r'^```(?:json)?\s*', '', content.strip())
                content = re.sub(r'\s*```$', '', content.strip())

                cleaned_content = AIGenerator._clean_json_string(content)

                json_match = re.search(r'\[.*\]', cleaned_content, re.DOTALL)
                if not json_match:
                    logging.error(f"OPENROUTER FALLBACK: AI response did not contain a JSON array. First 500 chars: {cleaned_content[:500]}")
                    return None

                raw_questions = json.loads(json_match.group())
                if not isinstance(raw_questions, list):
                    logging.error(f"OPENROUTER FALLBACK: Parsed JSON is not a list. Type: {type(raw_questions)}")
                    return None

                questions = AIGenerator._format_ai_questions(raw_questions, config_blocks)
                if not questions:
                    logging.error(f"OPENROUTER FALLBACK: AI returned {len(raw_questions)} questions, but ALL failed formatting/validation.")
                    return None

                print(f">>> [QUIZCRAFT DEBUG] SUCCESSFULLY PARSED {len(questions)} QUESTIONS <<<")
                return questions
            except requests.exceptions.HTTPError as e:
                status = e.response.status_code if e.response is not None else 0
                if status == 429:
                    print(f">>> [QUIZCRAFT DEBUG] HIT 429 RATE LIMIT. RETRYING IN 3 SECONDS... <<<")
                    time.sleep(3)
                    continue
                else:
                    error_body = ""
                    if e.response is not None:
                        try:
                            error_body = e.response.text
                        except:
                            pass
                    logging.error("="*50)
                    logging.error(f"OPENROUTER FALLBACK: NON-429 HTTP ERROR {status}")
                    logging.error(f"FULL ERROR BODY: {error_body}")
                    logging.error("="*50)
                    return None
            except json.JSONDecodeError as e:
                logging.error(f"OPENROUTER FALLBACK: AI returned invalid JSON. Error: {e}")
                return None
            except requests.exceptions.ConnectionError:
                logging.error("OPENROUTER FALLBACK: Connection failed - Could not reach OpenRouter API")
                return None
            except Exception as e:
                logging.error(f"OPENROUTER FALLBACK: Unexpected error - {str(e)}")
                return None

        logging.error("OPENROUTER FALLBACK: EXHAUSTED ALL 15 RETRIES DUE TO 429 RATE LIMITS.")
        return None

    @staticmethod
    def _format_ai_questions(raw_questions, config_blocks):
        type_queue = []
        for block in config_blocks:
            q_type = block.get('type', 'mcq')
            count = block.get('count', 5)
            marks = block.get('marks', 1)
            for _ in range(count):
                type_queue.append((q_type, marks))
        formatted = []
        for i, q in enumerate(raw_questions):
            if i < len(type_queue):
                expected_type, marks = type_queue[i]
            else:
                expected_type, marks = 'mcq', 1
            ai_type = str(q.get('type', 'mcq')).lower().strip()
            is_format_type = ai_type in ('mcq', 'msq', 'true_false')
            actual_type = ai_type if is_format_type else expected_type
            options = q.get('options', [])
            answer = q.get('answer', '')
            if actual_type == 'mcq':
                if not options or not isinstance(options, list) or len(options) < 2:
                    if isinstance(answer, str) and answer and len(answer) > 3:
                        options = [answer, "None of the above", "Partially correct", "Incorrect"]
                    else:
                        continue
                if isinstance(answer, list):
                    answer = answer[0] if answer else options[0]
                answer = str(answer).strip()
                if answer in ('A', 'B', 'C', 'D', 'E') and len(options) >= 4:
                    letter_to_idx = {'A': 0, 'B': 1, 'C': 2, 'D': 3, 'E': 4}
                    idx = letter_to_idx.get(answer, 0)
                    if idx < len(options):
                        answer = str(options[idx])
                if answer not in options:
                    options.insert(0, answer)
                options = [str(o) for o in options[:4]]
                unique_opts = []
                seen = set()
                for o in options:
                    if o.lower() not in seen:
                        seen.add(o.lower())
                        unique_opts.append(o)
                options = unique_opts
                if len(options) < 2:
                    continue
                if answer not in options:
                    answer = options[0]
            elif actual_type == 'msq':
                if not options or not isinstance(options, list):
                    continue
                options = [str(o) for o in options[:5]]
                unique_opts = []
                seen = set()
                for o in options:
                    if o.lower() not in seen:
                        seen.add(o.lower())
                        unique_opts.append(o)
                options = unique_opts
                if len(options) < 3:
                    continue
                if isinstance(answer, str):
                    answer_stripped = answer.strip()
                    if answer_stripped in ('A', 'B', 'C', 'D', 'E'):
                        answer = []
                    elif answer_stripped.startswith('[') and answer_stripped.endswith(']'):
                        try:
                            parsed = json.loads(answer_stripped)
                            if isinstance(parsed, list):
                                answer = [str(a) for a in parsed]
                            else:
                                answer = [answer_stripped]
                        except:
                            answer = [answer_stripped]
                    else:
                        answer = [answer_stripped]
                elif not isinstance(answer, list):
                    answer = []
                answer = [str(a).strip() for a in answer]
                valid_answer = []
                for a in answer:
                    if a in options:
                        valid_answer.append(a)
                    elif a in ('A', 'B', 'C', 'D', 'E'):
                        letter_to_idx = {'A': 0, 'B': 1, 'C': 2, 'D': 3, 'E': 4}
                        idx = letter_to_idx.get(a, -1)
                        if idx >= 0 and idx < len(options):
                            valid_answer.append(options[idx])
                answer = valid_answer
                if len(answer) < 2:
                    if len(options) >= 2:
                        answer = [options[0], options[1]]
                    else:
                        continue
            elif actual_type == 'true_false':
                answer = str(answer).strip()
                if answer not in ('True', 'False'):
                    if isinstance(answer, list) and len(answer) > 0:
                        first = str(answer[0]).strip()
                        if first in ('True', 'False'):
                            answer = first
                        else:
                            answer = 'True'
                    else:
                        answer = 'True'
                options = ['True', 'False']
            else:
                actual_type = expected_type
                if actual_type == 'mcq':
                    if not options or len(options) < 2:
                        continue
                    if isinstance(answer, list):
                        answer = answer[0]
                    answer = str(answer).strip()
                    if answer in ('A', 'B', 'C', 'D', 'E') and len(options) >= 4:
                        letter_to_idx = {'A': 0, 'B': 1, 'C': 2, 'D': 3, 'E': 4}
                        idx = letter_to_idx.get(answer, 0)
                        if idx < len(options):
                            answer = str(options[idx])
                    options = [str(o) for o in options[:4]]
                    if answer not in options:
                        options.insert(0, answer)
                elif actual_type == 'true_false':
                    answer = 'True'
                    options = ['True', 'False']
                else:
                    continue
            question_text = str(q.get('question', ''))
            if not question_text or len(question_text) < 15:
                continue
            formatted_q = {
                "question_text": question_text,
                "options": options,
                "answer": answer,
                "type": actual_type,
                "marks": marks,
            }
            if q.get('explanation'):
                formatted_q['explanation'] = str(q['explanation'])
            if q.get('concept'):
                formatted_q['concept'] = str(q['concept'])
            if q.get('thinking'):
                formatted_q['thinking_type'] = str(q['thinking'])
            formatted.append(formatted_q)
        return formatted


class QuizGenerator:

    def __init__(self):
        api_key = os.environ.get("OPENROUTER_API_KEY")
        if api_key:
            self.ai_available = AIGenerator.check_key_valid(api_key)
            print(f">>> [QUIZCRAFT DEBUG] AI MODE SET TO: {self.ai_available} <<<")
        else:
            self.ai_available = False

    def generate_from_text(self, text, config_blocks):
        if not text or not text.strip():
            return []
        text = text.strip()
        concepts = ConceptExtractor.extract(text)
        total_requested = sum(b.get('count', 0) for b in config_blocks)
        if self.ai_available:
            ai_questions = AIGenerator.generate_with_ai(text, concepts, config_blocks)
            if ai_questions and len(ai_questions) >= max(total_requested * 0.5, 1):
                return ai_questions
            if ai_questions:
                nlp_questions = NLPQuestionGenerator.generate(text, concepts, config_blocks)
                existing_texts = {q['question_text'].lower() for q in ai_questions}
                for nq in nlp_questions:
                    if nq['question_text'].lower() not in existing_texts:
                        ai_questions.append(nq)
                        if len(ai_questions) >= total_requested:
                            break
                return ai_questions[:total_requested]
        nlp_questions = NLPQuestionGenerator.generate(text, concepts, config_blocks)
        return nlp_questions